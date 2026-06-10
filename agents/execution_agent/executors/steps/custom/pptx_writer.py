"""
pptx_writer.py  —  PowerPoint Agent (Agent 06)
"""

import os
import re
from datetime import date
from ..base_step import BaseStep, StepResult
from ...core.envelope import resolve_path
from ...core import backend_client as bc

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def _rgb(hex_str: str) -> "RGBColor":
    """Convert a hex string like '1E2761' to RGBColor."""
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def _fill_bg(shape, hex_color: str) -> None:
    """Solid-fill a shape's background."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _set_font(run_or_para, font_name: str, size_pt: int,
              bold: bool = False, color_hex: str = None) -> None:
    font = run_or_para.font
    font.name = font_name
    font.size = Pt(size_pt)
    font.bold = bold
    if color_hex:
        font.color.rgb = _rgb(color_hex)


class PPTXWriter(BaseStep):

    # Default theme — used when LLM doesn't supply one
    _DEFAULT_THEME = {
        "primary":    "1E2761",   # navy
        "secondary":  "CADCFC",   # ice blue
        "accent":     "FFFFFF",   # white
        "title_font": "Georgia",
        "body_font":  "Calibri",
    }

    def run(self, envelope: dict, config: dict) -> StepResult:
        try:
            if not PPTX_AVAILABLE:
                return StepResult(
                    success=False, data={},
                    error="PPTXWriter: python-pptx is not installed."
                )

            output_dir     = config.get("output_dir", "output/presentations")
            naming_pattern = config.get("naming", "{task_id}_{date}_{short_title}.pptx")

            slide_data = resolve_path(
                envelope, "execution.steps.generate_slide_json.data"
            )

            if slide_data.get("paused_for_clarification", False):
                return StepResult(
                    success=True,
                    data={
                        "output_path": "", "slides_written": 0,
                        "template_used": "", "paused": True,
                        "clarification_question": slide_data.get(
                            "clarification_question", "Could you provide more details?"
                        ),
                        "status": "paused",
                    },
                    error=None
                )

            slides             = slide_data.get("slides", [])
            template_path      = slide_data.get("template_path", "")
            presentation_title = slide_data.get("presentation_title", "Presentation")
            theme              = {**self._DEFAULT_THEME, **slide_data.get("theme", {})}

            if not slides:
                return StepResult(
                    success=False, data={},
                    error="PPTXWriter: no slides found in generate_slide_json output."
                )

            if template_path and not template_path.startswith("templates/pptx/"):
                return StepResult(
                    success=False, data={},
                    error=f"PPTXWriter: template_path '{template_path}' is outside allowed directory."
                )

            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()
                # Widescreen 16:9
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)

            slides_written = 0
            for slide_spec in slides:
                self._write_slide(prs, slide_spec, theme)
                slides_written += 1

            task_id     = envelope.get("task", {}).get("task_id", "TASK-UNKNOWN")
            today       = date.today().strftime("%Y%m%d")
            short_title = self._slugify(presentation_title)[:30]
            filename    = (
                naming_pattern
                .replace("{task_id}",     task_id)
                .replace("{date}",        today)
                .replace("{short_title}", short_title)
            )

            if not output_dir.startswith("output/"):
                return StepResult(
                    success=False, data={},
                    error=f"PPTXWriter: output_dir '{output_dir}' is outside allowed sandbox."
                )

            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename)
            prs.save(output_path)

            # Upload to the task-outputs bucket + register on the task (best-effort).
            storage_path = bc.upload_and_register(output_path, envelope, "pptx")

            return StepResult(
                success=True,
                data={
                    "output_path":            output_path,
                    "storage_path":           storage_path,
                    "slides_written":         slides_written,
                    "template_used":          template_path or "blank",
                    "paused":                 False,
                    "clarification_question": "",
                },
                error=None
            )

        except KeyError as e:
            return StepResult(
                success=False, data={},
                error=f"PPTXWriter could not find required envelope path: {e}."
            )
        except Exception as e:
            return StepResult(
                success=False, data={},
                error=f"PPTXWriter unexpected error: {str(e)}"
            )

    # ── slide dispatcher ──────────────────────────────────────────────────────

    def _write_slide(self, prs, slide_spec: dict, theme: dict) -> None:
        layout = slide_spec.get("layout", "bullets")
        dispatch = {
            "title_slide":  self._layout_title_slide,
            "stat_callout": self._layout_stat_callout,
            "two_column":   self._layout_two_column,
            "bullets":      self._layout_bullets,
        }
        fn = dispatch.get(layout, self._layout_bullets)
        fn(prs, slide_spec, theme)

    # ── layouts ───────────────────────────────────────────────────────────────

    def _layout_title_slide(self, prs, spec: dict, theme: dict) -> None:
        """Dark full-bleed title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H  = prs.slide_width, prs.slide_height

        # Dark background
        bg = slide.shapes.add_shape(1, 0, 0, W, H)  # MSO_SHAPE_TYPE.RECTANGLE = 1
        _fill_bg(bg, theme["primary"])
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = spec.get("title", "")
        p.alignment = PP_ALIGN.LEFT
        _set_font(p, theme["title_font"], 44, bold=True, color_hex=theme["accent"])

        # Subtitle from first bullet
        bullets = spec.get("bullet_points", [])
        if bullets:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(10), Inches(1.2))
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = bullets[0]
            p2.alignment = PP_ALIGN.LEFT
            _set_font(p2, theme["body_font"], 18, color_hex=theme["secondary"])

        self._add_notes(slide, spec.get("speaker_notes", ""))

    def _layout_stat_callout(self, prs, spec: dict, theme: dict) -> None:
        """Slide with a large stat on the left, bullets on the right."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H  = prs.slide_width, prs.slide_height

        # Light background
        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_bg(bg, "F7F9FC")
        bg.line.fill.background()

        # Title bar
        self._add_title_bar(slide, spec.get("title", ""), theme, W)

        # Big stat block (left panel)
        stat = spec.get("stat") or {}
        stat_val   = stat.get("value", "")
        stat_label = stat.get("label", "")

        left_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(4.5), Inches(3.5))
        tf = left_box.text_frame
        tf.word_wrap = True

        p_val = tf.paragraphs[0]
        p_val.text = stat_val
        p_val.alignment = PP_ALIGN.CENTER
        _set_font(p_val, theme["title_font"], 64, bold=True, color_hex=theme["primary"])

        p_lbl = tf.add_paragraph()
        p_lbl.text = stat_label
        p_lbl.alignment = PP_ALIGN.CENTER
        _set_font(p_lbl, theme["body_font"], 16, color_hex="666666")

        # Bullet points (right panel)
        self._add_bullet_block(
            slide, spec.get("bullet_points", []),
            Inches(5.6), Inches(2.0), Inches(7.1), Inches(4.5), theme
        )

        self._add_notes(slide, spec.get("speaker_notes", ""))

    def _layout_two_column(self, prs, spec: dict, theme: dict) -> None:
        """Two equal bullet columns side by side."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H  = prs.slide_width, prs.slide_height

        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_bg(bg, "F7F9FC")
        bg.line.fill.background()

        self._add_title_bar(slide, spec.get("title", ""), theme, W)

        bullets = spec.get("bullet_points", [])
        mid     = (len(bullets) + 1) // 2
        left_bullets  = bullets[:mid]
        right_bullets = bullets[mid:]

        self._add_bullet_block(
            slide, left_bullets,
            Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), theme
        )
        self._add_bullet_block(
            slide, right_bullets,
            Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), theme
        )

        # Vertical divider
        line = slide.shapes.add_shape(1, Inches(6.5), Inches(2.0), Inches(0.04), Inches(4.2))
        _fill_bg(line, theme["secondary"])
        line.line.fill.background()

        self._add_notes(slide, spec.get("speaker_notes", ""))

    def _layout_bullets(self, prs, spec: dict, theme: dict) -> None:
        """Standard title + bullets slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H  = prs.slide_width, prs.slide_height

        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _fill_bg(bg, "F7F9FC")
        bg.line.fill.background()

        self._add_title_bar(slide, spec.get("title", ""), theme, W)

        self._add_bullet_block(
            slide, spec.get("bullet_points", []),
            Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.8), theme
        )

        self._add_notes(slide, spec.get("speaker_notes", ""))

    # ── shared helpers ────────────────────────────────────────────────────────

    def _add_title_bar(self, slide, title: str, theme: dict, W) -> None:
        """Colored left accent + title text at the top."""
        # Left accent bar
        bar = slide.shapes.add_shape(1, 0, Inches(0.9), Inches(0.08), Inches(0.85))
        _fill_bg(bar, theme["primary"])
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.25), W - Inches(0.5), Inches(1.0))
        tf = title_box.text_frame
        p  = tf.paragraphs[0]
        p.text = title
        _set_font(p, theme["title_font"], 28, bold=True, color_hex=theme["primary"])

    def _add_bullet_block(self, slide, bullets: list, left, top, width, height, theme: dict) -> None:
        """Adds a text box with bullet points."""
        if not bullets:
            return
        box = slide.shapes.add_textbox(left, top, width, height)
        tf  = box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            _set_font(p, theme["body_font"], 15, color_hex="1A1A2E")
            p.space_after = Pt(6)

    def _add_notes(self, slide, notes: str) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _slugify(self, text: str) -> str:
        text = text.lower()
        text = re.sub(r"[^a-z0-9\s]", "", text)
        text = re.sub(r"\s+", "_", text.strip())
        return text